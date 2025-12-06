import os
import time
import base64
import requests
from urllib.parse import unquote
from io import BytesIO
from supabase import create_client, Client
from openai import OpenAI
import tempfile

# Bibliotecas para procesamiento de documentos
import PyPDF2
import pdfplumber
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

# Bibliotecas para procesamiento de videos
import cv2
from PIL import Image

# ----------------------------------------------------
# 1. CONFIGURACIÓN E INICIALIZACIÓN
# ----------------------------------------------------

SUPABASE_URL = os.environ.get("SUPABASE_URL")
SUPABASE_SERVICE_KEY = os.environ.get("SUPABASE_SERVICE_KEY")
OPENAI_API_KEY = os.environ.get("OPENAI_API_KEY")
BUCKET_NAME = os.environ.get("SUPABASE_BUCKET", "whatsapp-media")

# Validar variables de entorno
if not all([SUPABASE_URL, SUPABASE_SERVICE_KEY, OPENAI_API_KEY]):
    raise ValueError("Faltan variables de entorno necesarias. Verifica SUPABASE_URL, SUPABASE_SERVICE_KEY y OPENAI_API_KEY")

# Inicializar clientes
supabase: Client = create_client(SUPABASE_URL, SUPABASE_SERVICE_KEY)
openai_client = OpenAI(api_key=OPENAI_API_KEY)

# Configuración de procesamiento de videos
VIDEO_FRAME_INTERVAL_SECONDS = 3  # Extraer un frame cada N segundos
VIDEO_MAX_FRAMES = 10  # Máximo de frames a analizar por video (para limitar costos)

# Mapeo de MIME types a extensiones
MIME_TYPE_MAP = {
    # Imágenes
    'image/jpeg': 'jpeg',
    'image/jpg': 'jpg',
    'image/png': 'png',
    'image/gif': 'gif',
    'image/webp': 'webp',
    # Documentos
    'application/pdf': 'pdf',
    'application/vnd.openxmlformats-officedocument.wordprocessingml.document': 'docx',
    'application/msword': 'doc',
    'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': 'xlsx',
    'application/vnd.ms-excel': 'xls',
    'application/vnd.openxmlformats-officedocument.presentationml.presentation': 'pptx',
    'application/vnd.ms-powerpoint': 'ppt',
}

SUPPORTED_IMAGE_FORMATS = ['png', 'jpeg', 'jpg', 'gif', 'webp']
SUPPORTED_DOCUMENT_FORMATS = ['pdf', 'docx', 'doc', 'xlsx', 'xls', 'pptx', 'ppt']
SUPPORTED_VIDEO_FORMATS = ['mp4', 'mov', 'avi', 'mkv', 'webm']

# ----------------------------------------------------
# 2. FUNCIONES DE UTILIDAD
# ----------------------------------------------------

def clean_url(url: str) -> str:
    """Limpia y decodifica la URL para evitar problemas de encoding."""
    try:
        return unquote(url)
    except Exception as e:
        print(f"⚠️ Error al limpiar URL: {e}")
        return url

def get_file_extension_from_url(url: str) -> str:
    """Extrae la extensión del archivo desde la URL."""
    url_lower = url.lower()
    
    # Intentar extraer extensión del final de la URL antes de parámetros
    # Ejemplo: file.pdf?token=xxx o file_lid.xlsx
    if '?' in url_lower:
        url_lower = url_lower.split('?')[0]
    
    # Buscar patrones comunes: _lid.ext, _false.ext, .ext
    import re
    
    # Patrón para encontrar extensiones comunes
    all_extensions = SUPPORTED_IMAGE_FORMATS + SUPPORTED_DOCUMENT_FORMATS + SUPPORTED_VIDEO_FORMATS
    pattern = r'[_\.](' + '|'.join(all_extensions) + r')(?:[_\?]|$)'
    match = re.search(pattern, url_lower)
    
    if match:
        return match.group(1)
    
    # Fallback: buscar extensión simple
    for ext in all_extensions:
        if f".{ext}" in url_lower:
            return ext
    
    return None

def get_file_metadata_from_storage(url: str) -> dict:
    """Obtiene metadata del archivo desde Supabase Storage."""
    try:
        # Extraer el path del archivo desde la URL
        # URL format: https://...supabase.co/storage/v1/object/public/BUCKET_NAME/path/to/file.ext
        if '/storage/v1/object/public/' in url:
            parts = url.split('/storage/v1/object/public/')
            if len(parts) > 1:
                # Remover el bucket name y obtener el path
                path_with_bucket = parts[1]
                path_parts = path_with_bucket.split('/', 1)
                if len(path_parts) > 1:
                    file_path = path_parts[1]
                    
                    # Limpiar encoding de la URL
                    file_path = unquote(file_path)
                    
                    # Consultar metadata desde Supabase Storage
                    bucket = supabase.storage.from_(BUCKET_NAME)
                    
                    # Listar archivos y buscar el correcto
                    # Nota: list() devuelve metadata incluyendo content_type
                    files = bucket.list()
                    
                    # Buscar el archivo específico
                    for file_info in files:
                        if file_info.get('name') in file_path:
                            return {
                                'content_type': file_info.get('metadata', {}).get('mimetype'),
                                'size': file_info.get('metadata', {}).get('size'),
                                'name': file_info.get('name')
                            }
        
        return None
        
    except Exception as e:
        print(f"⚠️ No se pudo obtener metadata del storage: {e}")
        return None

def get_file_extension(url: str, mime_type: str = None) -> str:
    """
    Determina la extensión del archivo.
    Prioridad: 1) MIME type, 2) Extensión de URL
    """
    # Intentar desde MIME type primero
    if mime_type and mime_type in MIME_TYPE_MAP:
        return MIME_TYPE_MAP[mime_type]
    
    # Fallback a extensión de URL
    return get_file_extension_from_url(url)

def download_file(url: str) -> bytes:
    """Descarga un archivo desde una URL y devuelve su contenido en bytes."""
    try:
        clean_file_url = clean_url(url)
        response = requests.get(clean_file_url, timeout=30)
        response.raise_for_status()
        return response.content
    except Exception as e:
        print(f"❌ Error al descargar archivo {url}: {e}")
        return None

# ----------------------------------------------------
# 3. PROCESAMIENTO DE IMÁGENES
# ----------------------------------------------------

def get_image_mime_type(extension: str) -> str:
    """Determina el tipo MIME de la imagen basado en la extensión."""
    mime_map = {
        'png': 'image/png',
        'jpg': 'image/jpeg',
        'jpeg': 'image/jpeg',
        'gif': 'image/gif',
        'webp': 'image/webp'
    }
    return mime_map.get(extension.lower(), 'image/jpeg')

def analyze_image_with_ai(image_base64: str, file_type: str) -> str:
    """Usa GPT-4o para obtener una descripción textual de la imagen."""
    
    prompt = (
        "Actúa como un analista experto de inteligencia de negocios en sector minero. "
        "Describe concisamente la imagen. Identifica cualquier texto relevante, "
        "avance de proyecto (si aplica), o problema visible. "
        "El objetivo es convertir la imagen en contexto textual para un reporte ejecutivo. "
        "Máximo 100 palabras."
    )

    try:
        response = openai_client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": prompt},
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:{file_type};base64,{image_base64}",
                            },
                        },
                    ],
                }
            ],
            max_tokens=200,
        )
        return response.choices[0].message.content.strip()

    except Exception as e:
        print(f"❌ Error en la API de OpenAI para la imagen: {e}")
        return None

def process_image(url: str, extension: str) -> str:
    """Procesa una imagen y retorna su análisis textual."""
    try:
        file_content = download_file(url)
        if not file_content:
            return None
        
        # Codificar a Base64
        image_base64 = base64.b64encode(file_content).decode('utf-8')
        
        # Analizar con IA
        file_type = get_image_mime_type(extension)
        description = analyze_image_with_ai(image_base64, file_type)
        
        return description
        
    except Exception as e:
        print(f"❌ Error procesando imagen: {e}")
        return None

# ----------------------------------------------------
# 4. PROCESAMIENTO DE PDFs
# ----------------------------------------------------

def extract_text_from_pdf(file_content: bytes) -> str:
    """Extrae texto de un archivo PDF."""
    text_parts = []
    
    try:
        # Método 1: PyPDF2 (más rápido)
        pdf_reader = PyPDF2.PdfReader(BytesIO(file_content))
        for page in pdf_reader.pages:
            page_text = page.extract_text()
            if page_text:
                text_parts.append(page_text)
        
        # Si PyPDF2 no extrajo texto, intentar con pdfplumber
        if not text_parts:
            with pdfplumber.open(BytesIO(file_content)) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(page_text)
        
        return "\n\n".join(text_parts)
        
    except Exception as e:
        print(f"❌ Error extrayendo texto de PDF: {e}")
        return None

def process_pdf(url: str) -> str:
    """Procesa un PDF y retorna su contenido textual."""
    try:
        file_content = download_file(url)
        if not file_content:
            return None
        
        text = extract_text_from_pdf(file_content)
        
        if text and len(text.strip()) > 0:
            return f"[CONTENIDO PDF]: {text[:3000]}"
        else:
            return "[PDF sin texto extraíble - posiblemente escaneado]"
            
    except Exception as e:
        print(f"❌ Error procesando PDF: {e}")
        return None

# ----------------------------------------------------
# 5. PROCESAMIENTO DE ARCHIVOS WORD
# ----------------------------------------------------

def extract_text_from_docx(file_content: bytes) -> str:
    """Extrae texto de un archivo Word (.docx)."""
    try:
        doc = Document(BytesIO(file_content))
        text_parts = []
        
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text)
        
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join([cell.text.strip() for cell in row.cells])
                if row_text.strip():
                    text_parts.append(row_text)
        
        return "\n".join(text_parts)
        
    except Exception as e:
        print(f"❌ Error extrayendo texto de DOCX: {e}")
        return None

def process_word(url: str) -> str:
    """Procesa un archivo Word y retorna su contenido textual."""
    try:
        file_content = download_file(url)
        if not file_content:
            return None
        
        text = extract_text_from_docx(file_content)
        
        if text and len(text.strip()) > 0:
            return f"[CONTENIDO WORD]: {text[:3000]}"
        else:
            return "[Documento Word vacío o sin contenido]"
            
    except Exception as e:
        print(f"❌ Error procesando Word: {e}")
        return None

# ----------------------------------------------------
# 6. PROCESAMIENTO DE ARCHIVOS EXCEL
# ----------------------------------------------------

def extract_text_from_xlsx(file_content: bytes) -> str:
    """Extrae texto de un archivo Excel (.xlsx)."""
    try:
        workbook = load_workbook(BytesIO(file_content), data_only=True)
        text_parts = []
        
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text_parts.append(f"\n=== HOJA: {sheet_name} ===")
            
            for row_idx, row in enumerate(sheet.iter_rows(values_only=True), 1):
                if row_idx > 100:
                    text_parts.append("[... contenido truncado ...]")
                    break
                
                row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                if row_text.strip():
                    text_parts.append(row_text)
        
        return "\n".join(text_parts)
        
    except Exception as e:
        print(f"❌ Error extrayendo texto de XLSX: {e}")
        return None

def process_excel(url: str) -> str:
    """Procesa un archivo Excel y retorna su contenido textual."""
    try:
        file_content = download_file(url)
        if not file_content:
            return None
        
        text = extract_text_from_xlsx(file_content)
        
        if text and len(text.strip()) > 0:
            return f"[CONTENIDO EXCEL]: {text[:3000]}"
        else:
            return "[Archivo Excel vacío o sin contenido]"
            
    except Exception as e:
        print(f"❌ Error procesando Excel: {e}")
        return None

# ----------------------------------------------------
# 7. PROCESAMIENTO DE ARCHIVOS POWERPOINT
# ----------------------------------------------------

def extract_text_from_pptx(file_content: bytes) -> str:
    """Extrae texto de un archivo PowerPoint (.pptx)."""
    try:
        presentation = Presentation(BytesIO(file_content))
        text_parts = []
        
        for slide_idx, slide in enumerate(presentation.slides, 1):
            text_parts.append(f"\n=== DIAPOSITIVA {slide_idx} ===")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text)
        
        return "\n".join(text_parts)
        
    except Exception as e:
        print(f"❌ Error extrayendo texto de PPTX: {e}")
        return None

def process_powerpoint(url: str) -> str:
    """Procesa un archivo PowerPoint y retorna su contenido textual."""
    try:
        file_content = download_file(url)
        if not file_content:
            return None
        
        text = extract_text_from_pptx(file_content)
        
        if text and len(text.strip()) > 0:
            return f"[CONTENIDO POWERPOINT]: {text[:3000]}"
        else:
            return "[Presentación vacía o sin contenido]"
            
    except Exception as e:
        print(f"❌ Error procesando PowerPoint: {e}")
        return None

# ----------------------------------------------------
# 7B. PROCESAMIENTO DE VIDEOS
# ----------------------------------------------------

def extract_frames_from_video(file_content: bytes, interval_seconds: int = 3, max_frames: int = 10) -> list:
    """
    Extrae frames de un video cada N segundos.
    
    Args:
        file_content: Contenido del video en bytes
        interval_seconds: Intervalo en segundos entre frames (default: 3)
        max_frames: Máximo número de frames a extraer (default: 10)
    
    Returns:
        Lista de imágenes PIL
    """
    try:
        # Guardar temporalmente el video
        with tempfile.NamedTemporaryFile(suffix='.mp4', delete=False) as tmp_file:
            tmp_file.write(file_content)
            tmp_path = tmp_file.name
        
        # Abrir video con OpenCV
        video = cv2.VideoCapture(tmp_path)
        
        # Obtener información del video
        total_frames = int(video.get(cv2.CAP_PROP_FRAME_COUNT))
        fps = video.get(cv2.CAP_PROP_FPS)
        duration = total_frames / fps if fps > 0 else 0
        
        print(f"      📹 Video: {duration:.1f}s, FPS: {fps:.1f}, {total_frames} frames totales")
        
        frames = []
        
        if total_frames > 0 and fps > 0:
            # Calcular frames a extraer cada N segundos
            frames_per_interval = int(fps * interval_seconds)
            
            # Calcular posiciones de frames
            frame_positions = []
            current_frame = 0
            
            while current_frame < total_frames and len(frame_positions) < max_frames:
                frame_positions.append(current_frame)
                current_frame += frames_per_interval
            
            print(f"      🎞️ Extrayendo {len(frame_positions)} frames (cada {interval_seconds}s)...")
            
            # Extraer frames
            for idx, frame_pos in enumerate(frame_positions):
                video.set(cv2.CAP_PROP_POS_FRAMES, frame_pos)
                ret, frame = video.read()
                
                if ret:
                    # Convertir de BGR (OpenCV) a RGB (PIL)
                    frame_rgb = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
                    pil_image = Image.fromarray(frame_rgb)
                    frames.append(pil_image)
                    
                    # Calcular timestamp para referencia
                    timestamp = frame_pos / fps
                    print(f"         ✓ Frame {idx+1} extraído (t={timestamp:.1f}s)")
        
        video.release()
        
        # Eliminar archivo temporal
        try:
            os.unlink(tmp_path)
        except:
            pass
        
        return frames
        
    except Exception as e:
        print(f"❌ Error extrayendo frames del video: {e}")
        return []

def analyze_video_frame(frame_image: Image, frame_number: int, timestamp: float) -> str:
    """Analiza un frame de video usando GPT-4o."""
    
    prompt = (
        f"Analiza este frame de un video de WhatsApp relacionado con operaciones mineras (timestamp: {timestamp:.1f}s). "
        "Describe lo que ves: equipos, operaciones, personas, problemas, condiciones, o cualquier elemento relevante. "
        "Si hay texto visible (pantallas, letreros, medidores), transcríbelo. "
        "Si identificas un problema o situación de riesgo, menciónalo. "
        "Sé conciso y específico. Máximo 100 palabras."
    )
    
    try:
        # Convertir PIL Image a base64
        buffer = BytesIO()
        frame_image.save(buffer, format='JPEG', quality=85)
        image_base64 = base64.b64encode(buffer.getvalue()).decode('utf-8')
        
        response = openai_client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": prompt},
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/jpeg;base64,{image_base64}",
                            },
                        },
                    ],
                }
            ],
            max_tokens=200,
        )
        return response.choices[0].message.content.strip()
        
    except Exception as e:
        print(f"❌ Error analizando frame con IA: {e}")
        return None

def process_video(url: str) -> str:
    """
    Procesa un video extrayendo y analizando frames según configuración.
    Usa VIDEO_FRAME_INTERVAL_SECONDS y VIDEO_MAX_FRAMES definidos globalmente.
    """
    try:
        print(f"      🎬 Descargando video...")
        file_content = download_file(url)
        if not file_content:
            return None
        
        # Extraer frames según configuración
        frames = extract_frames_from_video(
            file_content, 
            interval_seconds=VIDEO_FRAME_INTERVAL_SECONDS,
            max_frames=VIDEO_MAX_FRAMES
        )
        
        if not frames:
            return "[Video sin frames extraíbles]"
        
        print(f"      🤖 Analizando {len(frames)} frames con IA...")
        
        # Analizar cada frame
        frame_analyses = []
        for idx, frame in enumerate(frames):
            timestamp = idx * VIDEO_FRAME_INTERVAL_SECONDS
            analysis = analyze_video_frame(frame, idx + 1, timestamp)
            if analysis:
                frame_analyses.append(f"[t={timestamp}s] {analysis}")
                print(f"         ✓ Frame {idx+1} analizado")
        
        if frame_analyses:
            combined_analysis = "\n".join(frame_analyses)
            return f"[ANÁLISIS DE VIDEO - {len(frames)} frames cada {VIDEO_FRAME_INTERVAL_SECONDS}s]:\n{combined_analysis}"
        else:
            return "[Video procesado pero sin análisis disponible]"
            
    except Exception as e:
        print(f"❌ Error procesando video: {e}")
        import traceback
        traceback.print_exc()
        return None

# ----------------------------------------------------
# 8. PROCESADOR UNIVERSAL DE ARCHIVOS
# ----------------------------------------------------

def process_file(url: str, file_extension: str) -> str:
    """Procesa cualquier tipo de archivo soportado y retorna su contenido."""
    
    print(f"   📄 Tipo de archivo detectado: .{file_extension}")
    
    # Imágenes
    if file_extension in SUPPORTED_IMAGE_FORMATS:
        return process_image(url, file_extension)
    
    # Videos
    elif file_extension in SUPPORTED_VIDEO_FORMATS:
        return process_video(url)
    
    # PDFs
    elif file_extension == 'pdf':
        return process_pdf(url)
    
    # Word
    elif file_extension in ['docx', 'doc']:
        return process_word(url)
    
    # Excel
    elif file_extension in ['xlsx', 'xls']:
        return process_excel(url)
    
    # PowerPoint
    elif file_extension in ['pptx', 'ppt']:
        return process_powerpoint(url)
    
    else:
        print(f"   ⚠️ Formato no soportado: .{file_extension}")
        return None

# ----------------------------------------------------
# 9. GENERACIÓN DE EMBEDDINGS
# ----------------------------------------------------

def create_and_upload_embedding(content: str, record_id: int):
    """Genera el embedding y actualiza el registro en Supabase."""
    
    try:
        # 1. Generar Embedding
        print(f"   [ID {record_id}] Generando embedding...")
        embedding_response = openai_client.embeddings.create(
            input=content[:8000],  # Limitar a 8000 caracteres para evitar límites de tokens
            model="text-embedding-3-small"
        )
        embedding_vector = embedding_response.data[0].embedding

        # 2. Actualizar Supabase
        update_response = supabase.from_('mensajes_analisis').update({
            'embedding': embedding_vector,
            'procesado_ia': True
        }).eq('id', record_id).execute()

        if update_response.data:
            print(f"✔️ Actualizado ID {record_id} con embedding.")
            return True
        else:
            print(f"❌ Error al actualizar ID {record_id}.")
            return False
            
    except Exception as e:
        print(f"❌ Error procesando ID {record_id}: {e}")
        return False

# ----------------------------------------------------
# 10. LÓGICA PRINCIPAL
# ----------------------------------------------------

def main_processor():
    """Procesa mensajes pendientes de vectorización."""
    print("\n" + "="*70)
    print("🚀 Iniciando Proceso de Vectorización y Análisis de Documentos")
    print("="*70)

    try:
        # 1. Buscar registros sin vectorizar
        query_response = supabase.from_('mensajes_analisis').select("*").is_('embedding', 'null').order('fecha_hora', desc=False).limit(50).execute()
        
        pending_records = query_response.data if query_response.data else []

        if not pending_records:
            print("✅ No hay nuevos registros para procesar.")
            return

        print(f"🔎 Encontrados {len(pending_records)} registros pendientes.")

        # 2. Procesar cada registro
        processed_count = 0
        error_count = 0
        skipped_count = 0
        
        for idx, record in enumerate(pending_records, 1):
            record_id = record.get('id')
            final_content = record.get('contenido_texto', '') or ""

            print(f"\n{'─'*70}")
            print(f"📝 [{idx}/{len(pending_records)}] Procesando ID {record_id}...")

            # A. Si tiene archivo adjunto
            if record.get('url_storage'):
                file_url = record['url_storage']
                
                # Intentar obtener metadata del storage
                # metadata = get_file_metadata_from_storage(file_url)
                # mime_type = metadata.get('content_type') if metadata else None
                
                # Por ahora usar extensión de URL
                file_extension = get_file_extension_from_url(file_url)
                
                if file_extension:
                    print(f"   [ID {record_id}] 📎 Procesando archivo adjunto...")
                    
                    # Procesar archivo según su tipo
                    file_content = process_file(file_url, file_extension)
                    
                    if file_content:
                        final_content = f"{final_content}\n\n{file_content}"
                        print(f"   [ID {record_id}] ✅ Archivo procesado exitosamente")
                    else:
                        print(f"   [ID {record_id}] ⚠️ No se pudo procesar el archivo")
                else:
                    print(f"   [ID {record_id}] ⚠️ Tipo de archivo no reconocido")
                    print(f"   [ID {record_id}] 🔗 URL: {file_url[:100]}...")  # Mostrar primeros 100 caracteres

            # B. Vectorización del contenido final
            if final_content.strip():
                success = create_and_upload_embedding(final_content, record_id)
                if success:
                    processed_count += 1
                else:
                    error_count += 1
            else:
                print(f"   [ID {record_id}] ⚠️ Contenido vacío. Saltando.")
                skipped_count += 1

        print(f"\n{'='*70}")
        print(f"📊 RESUMEN DEL CICLO:")
        print(f"   ✅ Procesados: {processed_count}")
        print(f"   ❌ Errores: {error_count}")
        print(f"   ⏭️  Saltados: {skipped_count}")
        print(f"{'='*70}")

    except Exception as e:
        print(f"❌ Error en main_processor: {e}")
        import traceback
        traceback.print_exc()

# ----------------------------------------------------
# 11. PUNTO DE ENTRADA
# ----------------------------------------------------

if __name__ == "__main__":
    print("="*70)
    print("🔧 Servicio de Vectorización Multi-Formato")
    print("="*70)
    print(f"🌐 Conectado a Supabase: {SUPABASE_URL}")
    print(f"📁 Bucket: {BUCKET_NAME}")
    print(f"📄 Formatos soportados:")
    print(f"   • Imágenes: {', '.join(SUPPORTED_IMAGE_FORMATS)}")
    print(f"   • Videos: {', '.join(SUPPORTED_VIDEO_FORMATS)}")
    print(f"     └─ Configuración: 1 frame cada {VIDEO_FRAME_INTERVAL_SECONDS}s (máx {VIDEO_MAX_FRAMES} frames)")
    print(f"   • Documentos: {', '.join(SUPPORTED_DOCUMENT_FORMATS)}")
    print("="*70)
    print("🔍 Búsqueda Semántica RAG: HABILITADA ✅")
    print("⏰ El servicio verifica nuevos registros cada 30 segundos")
    print("🔄 Para detener el servicio, presiona Ctrl+C")
    print("="*70 + "\n")
    
    # Bucle continuo para servicio 24/7
    cycle_count = 0
    while True:
        try:
            cycle_count += 1
            timestamp = time.strftime('%Y-%m-%d %H:%M:%S')
            print(f"\n🔄 Ciclo #{cycle_count} - {timestamp}")
            
            main_processor()
            
            print(f"\n😴 Esperando 30 segundos antes del siguiente ciclo...")
            time.sleep(30)
            
        except KeyboardInterrupt:
            print("\n\n" + "="*70)
            print("👋 Servicio detenido por el usuario")
            print("="*70)
            break
        except Exception as e:
            print(f"\n❌ Error crítico en el ciclo principal: {e}")
            import traceback
            traceback.print_exc()
            print("⏰ Esperando 60 segundos antes de reintentar...")
            time.sleep(60)
